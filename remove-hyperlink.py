from pptx import Presentation

def remove_hyperlinks(file_path, output_path):
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.hyperlink.address = None
    prs.save(output_path)

input_file = "zjubeamer.pptx"  # 输入文件路径
output_file = "output.pptx"       # 输出文件路径
remove_hyperlinks(input_file, output_file)

